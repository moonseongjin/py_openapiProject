# ===============================================================================================================
# ===============================================================================================================
# 1. 종목코드 및 일별 시세 가져오기
import pandas as pd
import requests

# 함수정의

def get_stock_code():
    # 종목코드 다운로드
    stock_code = pd.read_html('http://kind.krx.co.kr/corpgeneral/corpList.do?method=download', header=0)[0]
    
    # 필요없는 column들은 제외
    stock_code = stock_code[['회사명', '종목코드']]
    
    # 한글 컬럼명을 영어로 변경
    stock_code = stock_code.rename(columns={'회사명': 'company', '종목코드': 'code'})
    
    # 종목코드가 6자리이기 때문에 6자리를 맞춰주기 위해 설정해줌
    stock_code.code = stock_code.code.map('{:06d}'.format)
    
    return stock_code

def get_stock(code):
    df_list = []
    for page in range(1, 21):
        # 일별 시세 url
        url = 'https://finance.naver.com/item/sise_day.nhn?code={code}&page={page}'.format(code=code, page=page)

        header = {'User-Agent': 'ㅇㅇㅇ'}

        res = requests.get(url, headers=header)
        
        # 테이블이 페이지에 존재하는지 확인
        tables = pd.read_html(res.text, header=0)
        if not tables:
            print(f"{page}페이지에서 테이블을 찾을 수 없습니다. 데이터 수집 중단.")
            break
        
        try:
            current_df = tables[0]
            df_list.append(current_df)
        except Exception as e:
            print(f"Error occurred on page {page}: {e}")

        df = pd.concat(df_list, ignore_index=True)
    return df

def clean_data(df):
    # df.dropna()를 이용해 결측값 있는 행 제거
    df = df.dropna()
    
    # 한글로 된 컬럼명을 영어로 바꿔줌
    df = df.rename(columns= {'날짜': 'date', '종가':'close','전일비':'diff','시가':'open', '고가':'high', '저가':'low','거래량':'volume'})
    
    # 데이터의 타입을 int형으로 바꿔줌
    df[['close','diff','open','high','low','volume']]=df[['close','diff','open','high','low','volume']].astype(int)
    
    # 컬럼명 'date'의 타입을 date로 바꿔줌
    df['date'] = pd.to_datetime(df['date'])
    
    # 일자(date)를 기준으로 오름차순 정렬
    df = df.sort_values(by = ['date'], ascending=True)
    
    return df

# 함수 호출

# 종목코드 가져오기
company='삼성전자'
stock_code = get_stock_code()

# 일별 시세가져오기
code = stock_code[stock_code.company==company].code.values[0].strip()
# strip() : 공백 제거
df = get_stock(code)

# 일별 시세 클린징
df = clean_data(df)
print(df)

# ===============================================================================================================
# ===============================================================================================================
# 2. 보고자료 준비하기
import matplotlib.pyplot as plt
from pandas.plotting import table
import os

plt.figure(figsize=(10,4))
plt.plot(df['date'], df['close'])
plt.xlabel('date')
plt.ylabel('close')

# 차트 저장 및 출력
chart_fname = os.path.join("res/stock_report", '{company}_chart.png'.format(company=company))
plt.savefig(chart_fname)
plt.show()

# 일별 시세 그리기
plt.figure(figsize=(15,4))
ax = plt.subplot(111,frame_on=False)
ax.xaxis.set_visible(False)
ax.yaxis.set_visible(False)
df = df.sort_values(by=['date'], ascending=False)
table(ax, df.head(10), loc='center', cellLoc='center', rowLoc = 'center')

# 일별 시세 저장하기
table_fname = os.path.join("res/stock_report", '{company}_table.png'.format(company=company))
plt.savefig(table_fname)

# ===============================================================================================================
# ===============================================================================================================
# 3. 보고서 작성하기

# 선행 pip install python-pptx

import datetime
from pptx import Presentation 
from pptx.util import Inches
import os

# 파워포인트 객체 선언
today = datetime.datetime.today().strftime('%Y%m%d')
prs = Presentation() # 파워포인트 객체 선언

# 제목 슬라이드 추가
title_slide_layout = prs.slide_layouts[0] # 0: 제목 슬라이드
slide = prs.slides.add_slide(title_slide_layout) # 재목 슬라이드를 파워포인트 객체에 추가

# 제목 - 제목에 값넣기
title = slide.shapes.title # 제목
title.text = "주식 보고서" # 제목에 값 넣기

# 부제목
subtitle = slide.placeholders[1] # 제목 상자는 placeholders[0], 부제목 상자는 [1]
subtitle.text = "보고서 작성일 : {date}".format(date=today)

# 차트 및 테이블 슬라이드 추가
title_only_slide_layout = prs.slide_layouts[5]
slide = prs.slides.add_slide(title_only_slide_layout)

shapes = slide.shapes
shapes.title.text = '{company}, {close}원에 거래 마감'.format(company=company, close=df.iloc[0]['close'])
print(shapes.title.text)
# 차트 추가

left = Inches(0.5)
height = Inches(2.5)
width = Inches(9)
top = Inches(2)
# width, hegith가 없을 경우 원본 사이즈로
pic = slide.shapes.add_picture(chart_fname, left, top, width=width, height=height)

# 테이블 추가
left = Inches(-1)
height = Inches(3)
width = Inches(12)
top = Inches(4)

pic = slide.shapes.add_picture(table_fname, left,top, width=width, height=height)
cursor_sp=slide.shapes[0].element
cursor_sp.addprevious(pic._element) # 해당 요소를 뒤로 보냅니다.

# 보고서 작성
ppt_fname = os.path.join("res/stock_report", 'stock_report.pptx')
prs.save(ppt_fname)

# ===============================================================================================================
# ===============================================================================================================
# 4. 보고서를 이메일로 전송하기
import smtplib

# 이메일 메시지를 이진 데이터로 바꿔주는 인코더
from email import encoders

# 텍스트 형식
from email.mime.text import MIMEText

# 이미지 형식
from email.mime.image import MIMEImage

# 오디오 형식
from email.mime.audio import MIMEAudio

# 이메일 메시지에 다양한 형식을 중첩하여 담기 위한 객체
from email.mime.multipart import MIMEMultipart

# 위의 모든 객체들을 생성할 수 있는 기본 객체
# MIMEBase(_maintype, _subtype)
# MIMEbASE(<메인 타입>, <서브 타입>)
from email.mime.base import MIMEBase

# 함수 정의

# 메일 전송
def send_email(send_info, msg):
    
    # SMTP 세션 생성
    with smtplib.SMTP(send_info["send_server"], send_info["send_port"]) as server:
        
        # TLS 보안 연결
        server.starttls()
        
        # 보내는사람 계정으로 로그인
        server.login(send_info["send_user_id"], send_info["send_user_pw"])

        # 로그인 된 서버에 이메일 전송
        response = server.sendmail(msg['from'], msg['to'], msg.as_string())

        # 이메일 전송 성공시
        if not response:
            print('이메일을 정상적으로 보냈습니다.')
        else:
            print(response)

# 첨부파일 추가
def attach_multi(multi_info):

    # 3. 이메일 메시지에 다양한 형식을 담아내기 위한 객체 설정
    multi = MIMEMultipart(_subtype='mixed')
    
    # 4. 이미지 갯수만큼 반복
    for key, value in multi_info.items():
        if key == 'text':
            with open(value['filename'], 'rb') as fp:
                msg = MIMEText(fp.read(), _subtype=value['subtype'])
        elif key == 'image':
            with open(value['filename'], 'rb') as fp:
                msg = MIMEImage(fp.read(), _subtype=value['subtype'])
        elif key == 'audio':
            with open(value['filename'], 'rb') as fp:
                msg = MIMEAudio(fp.read(), _subtype=value['subtype'])
        else:
            with open(value['filename'], 'rb') as fp:
                msg = MIMEBase(value['maintype'], value['subtype'])
                msg.set_payload(fp.read())
                encoders.encode_base64(msg)      
        
        # 파일명을 첨부파일 제목으로 추가
        
        # os.path.split() 함수는 파일 경로를 디렉토리 부분과 파일명 부분으로 나누어주는 함수
        # dirname에는 C:/Users/moonw/py_life/res/stock_report가, fname에는 stock_report.pptx
        dirname, fname = os.path.split(value['filename'])
        
        print(fname)  # 첨부 파일의 이름을 출력합니다.
        
        msg.add_header('Content-Disposition', 'attachment', filename=fname)
        multi.attach(msg)
    
    return multi

# 함수 호출
                      
# 메일 전송 파라미터
send_info = dict(
    {"send_server" : "smtp.naver.com", # SMTP서버 주소
     "send_port" : 587, # SMTP서버 포트
     "send_user_id" : "<송신자ID>@naver.com",
     "send_user_pw" : "비밀번호"
    })

multi_info = {
    'application': {'maintype': 'application', 'subtype': 'octet-stream', 'filename': 'C:/Users/moonw/py_life/res/stock_report/stock_report.pptx'},
    'image': {'maintype': 'image', 'subtype': 'png', 'filename': 'C:/Users/moonw/py_life/res/stock_report/삼성전자_chart.png'},
    'image2': {'maintype': 'image', 'subtype': 'png', 'filename': 'C:/Users/moonw/py_life/res/stock_report/삼성전자_table.png'}
}

title = '({date}).주식 보고서 분석 자료입니다.'.format(date=today)
content = "주식 보고서 분석 자료입니다."
sender = send_info["send_user_id"]
receiver = "@naver.com"
msg = MIMEText(_text = content, _charset = "utf-8")

# 첨부파일 추가
multi_info['application']['filename'] = ppt_fname

# 첨부파일 추가한 객체에 정보 세팅
multi = attach_multi(multi_info)

multi['subject'] = title # 메일 제목
multi['from'] = sender # 보낸사람
multi['to'] = receiver # 받는사람
multi.attach(msg)

# 메일 전송 함수를 호출
send_email(send_info, multi)