""" 
파워포인트 편집을 위해 필요한 라이브러리

pip install python-pptx

ppt 열어놓은 채로 실행하면 에러남.
"""

# ptt 생성 및 저장
from pptx import Presentation # 라이브러리
from pptx.util import Inches # 사진, 표 등을 그리기 위해

prs = Presentation() # 파워포인트 객체 선언

for i in range(0,11): 
    # 0은 제목 슬라이드
    # 1은 제목 밑 내용 슬라이드
    # 2는 구역 머리글 슬라이드 ... 
    title_slide_layout = prs.slide_layouts[i] # 슬라이드 종류 선택
    slide = prs.slides.add_slide(title_slide_layout) # 슬라이드 추가

# 저장할 ptt 파일이름
prs.save('add all slides.pptx')



