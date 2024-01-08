# 제목 및 내용 슬라이드

from pptx import Presentation # 라이브러리
from pptx.util import Inches # 사진, 표 등을 그리기 위해

prs = Presentation() # 파워포인트 객체 선언

two_layout = prs.slide_layouts[1] # 0은 제목 슬라이드
slide = prs.slides.add_slide(two_layout) # 슬라이드 추가

# 제목 - 제목에 값 넣기
title = slide.placeholders[0] # 제목
title.text = "Bullet 슬라이드 추가" # 제목에 값 넣기

# 내용
body_shape = slide.placeholders[1]
tf = body_shape.text_frame
tf.text = 'Find the bullet slide layout'

# 단락 추가
p = tf.add_paragraph()
p.text = 'hahahoho'
p.level = 1 # 1: 들여쓰기 레벨

# 단락 추가
p = tf.add_paragraph()
p.text = 'I am 산타.'
p.level = 2 

# 저장 
prs.save('test.pptx')