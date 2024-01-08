# 플레이스홀더에 대한 정보 출력
from pptx import Presentation # 라이브러리
from pptx.util import Inches # 사진, 표 등을 그리기 위해

prs = Presentation() # 파워포인트 객체 선언

for i in range(0,11): 
    print("[%d]"%(i))
    slide = prs.slides.add_slide(title_slide_layout) # 슬라이드 추가
    for shape in slide.placeholders:
        print('%d %s' % (shape.placeholder_format.idx, shape.name))